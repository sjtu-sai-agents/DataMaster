#!/usr/bin/env python3
"""Create an editable schematic walkthrough figure for DataMaster.

This is a paper-style schematic walkthrough for the random-acts-of-pizza task.
The script only reads existing run logs to sanity-check semi-concrete labels; it
never runs experiments, training, evaluation, grading, or submission code.

Primary output is an editable PPTX built from native PowerPoint shapes. PDF/PNG
are exported with LibreOffice when available, otherwise generated from a matched
matplotlib fallback layout.
"""

from __future__ import annotations

import json
import math
import shutil
import subprocess
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

import matplotlib.pyplot as plt
from matplotlib.patches import Circle, FancyArrowPatch, FancyBboxPatch, Polygon

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - fallback handled in main()
    Presentation = None


REPO_ROOT = Path(__file__).resolve().parents[1]
RUN_DIR = REPO_ROOT / "runs" / "ml_master_datatree_20260428_001912"
OUT_DIR = REPO_ROOT / "scripts" / "extraction" / "datamaster_walkthrough"
PPTX_PATH = OUT_DIR / "datamaster_random_acts_walkthrough.pptx"
PDF_PATH = OUT_DIR / "datamaster_random_acts_walkthrough.pdf"
PNG_PATH = OUT_DIR / "datamaster_random_acts_walkthrough.png"

SLIDE_W = 13.33
SLIDE_H = 7.50
FONT = "Arial"
MPL_FONT = "DejaVu Sans"

COLORS = {
    "white": "FFFFFF",
    "ink": "1F2937",
    "muted": "69758A",
    "faint": "E7ECF3",
    "line": "CBD3DF",
    "line_dark": "6B7280",
    "task_fill": "FFF3C4",
    "task_line": "D6A514",
    "init_fill": "CBD8E6",
    "init_line": "62748B",
    "red_fill": "F8D7D4",
    "red_line": "B7565B",
    "black_fill": "222A36",
    "black_line": "0F172A",
    "side_red": "F7E2E1",
    "side_black": "A7ADB6",
    "memory_fill": "E7F3E9",
    "memory_line": "5FA274",
    "memory_text": "2E6F46",
    "output_fill": "E6F0FA",
    "output_line": "5B82B8",
    "gold": "D39B05",
    "callout_fill": "FFFFFF",
    "callout_line": "D9E0EA",
}

# Layout constants in slide inches. These are intentionally centralized so the
# figure can be tuned without touching drawing logic.
NODE_D = 0.48
BEST_D = 0.58
SMALL_D = 0.26

TREE = {
    "task": {"kind": "box", "x": 5.05, "y": 0.34, "w": 2.50, "h": 0.60},
    "init": {"kind": "node", "cx": 6.30, "cy": 1.18, "d": NODE_D},
    "R1": {"kind": "node", "cx": 5.55, "cy": 1.86, "d": NODE_D},
    "B1": {"kind": "node", "cx": 7.05, "cy": 1.86, "d": NODE_D},
    "R2": {"kind": "node", "cx": 5.55, "cy": 2.56, "d": NODE_D},
    "B2": {"kind": "node", "cx": 7.05, "cy": 2.56, "d": NODE_D},
    "R3": {"kind": "node", "cx": 5.55, "cy": 3.26, "d": NODE_D},
    "B3": {"kind": "node", "cx": 7.05, "cy": 3.26, "d": NODE_D},
    "cycle": {"kind": "box", "x": 5.02, "y": 3.86, "w": 2.56, "h": 0.40},
    "Rk": {"kind": "node", "cx": 5.55, "cy": 4.58, "d": NODE_D},
    "Bk": {"kind": "node", "cx": 7.05, "cy": 4.58, "d": NODE_D},
    "Best": {"kind": "node", "cx": 6.30, "cy": 5.36, "d": BEST_D},
    "output": {"kind": "box", "x": 5.06, "y": 6.13, "w": 2.48, "h": 0.68},
}

MAIN_EDGES = [
    ("task", "init"),
    ("init", "R1"),
    ("R1", "B1"),
    ("B1", "R2"),
    ("R2", "B2"),
    ("B2", "R3"),
    ("R3", "B3"),
    ("B3", "cycle"),
    ("cycle", "Rk"),
    ("Rk", "Bk"),
    ("Bk", "Best"),
    ("Best", "output"),
]

STEP_CALLOUTS = [
    (0, "Initialize pizza-request task.", 0.45, 0.45, 3.45, 0.42, "task", "left", "task_line"),
    (1, "Inspect request text and metadata.", 0.45, 1.64, 3.45, 0.42, "R1", "left", "red_line"),
    (3, "Search sentiment/politeness signals.", 0.45, 2.34, 3.45, 0.42, "R2", "left", "red_line"),
    (5, "Explore requester/subreddit behavior.", 0.45, 3.04, 3.45, 0.42, "R3", "left", "red_line"),
    (7, "Read memory; refine branch.", 0.45, 4.37, 3.45, 0.42, "Rk", "left", "memory_line"),
    (9, "Select best leaf node.", 0.45, 5.15, 3.45, 0.42, "Best", "left", "gold"),
    (2, "Clean fields; train baseline variant.", 9.02, 1.64, 3.55, 0.42, "B1", "right", "black_fill"),
    (4, "Implement lexical feature pipeline.", 9.02, 2.34, 3.55, 0.42, "B2", "right", "black_fill"),
    (6, "Build user-history features; validate.", 9.02, 3.04, 3.55, 0.42, "B3", "right", "black_fill"),
    (8, "Repeated RED/BLACK refinement cycles.", 9.02, 3.86, 3.55, 0.42, "cycle", "right", "line_dark"),
    (10, "Produce final submission.", 9.02, 6.23, 3.55, 0.42, "output", "right", "output_line"),
]

MEMORY_BOX = {"x": 9.10, "y": 4.58, "w": 3.20, "h": 1.34}
ROLE_BOX = {"x": 9.06, "y": 0.42, "w": 3.46, "h": 0.88}

FALLBACK_LABELS = [label for _, label, *_ in STEP_CALLOUTS]


def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    hex_color = hex_color.strip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def ppt_rgb(hex_color: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(hex_color))


def load_json(path: Path) -> Optional[Any]:
    try:
        with path.open("r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return None


def iter_dicts(obj: Any) -> Iterable[Dict[str, Any]]:
    if isinstance(obj, dict):
        yield obj
        for value in obj.values():
            yield from iter_dicts(value)
    elif isinstance(obj, list):
        for item in obj:
            yield from iter_dicts(item)


def parse_run_logs() -> Dict[str, Any]:
    """Read existing logs only to confirm the schematic labels are grounded."""
    meta: Dict[str, Any] = {
        "ok": False,
        "warnings": [],
        "node_count": 0,
        "trajectory_records": 0,
        "best_node": None,
        "best_score": None,
        "max_depth": None,
        "keywords": {},
    }

    uct_path = RUN_DIR / "logs" / "uct_nodes" / "node.json"
    traj_path = RUN_DIR / "trajectories" / "task_0" / "trajectory.json"
    grade_path = RUN_DIR / "trajectories" / "task_0" / "grade_results.json"

    node_data = load_json(uct_path)
    if node_data is None:
        meta["warnings"].append(f"could not parse {uct_path}")
    else:
        nodes = []
        depths = []
        for d in iter_dicts(node_data):
            keys = set(d)
            if {"id", "stage"} & keys or {"node_id", "stage"} & keys:
                nodes.append(d)
                depth = d.get("depth") or d.get("node_depth")
                if isinstance(depth, (int, float)):
                    depths.append(int(depth))
        meta["node_count"] = len(nodes)
        if depths:
            meta["max_depth"] = max(depths)

    traj_data = load_json(traj_path)
    if traj_data is None:
        meta["warnings"].append(f"could not parse {traj_path}")
    else:
        if isinstance(traj_data, list):
            meta["trajectory_records"] = len(traj_data)
        elif isinstance(traj_data, dict):
            records = traj_data.get("trajectory") or traj_data.get("records") or traj_data.get("messages")
            meta["trajectory_records"] = len(records) if isinstance(records, list) else 1
        small_text = json.dumps(traj_data, ensure_ascii=True).lower()
        for key in [
            "request",
            "metadata",
            "sentiment",
            "politeness",
            "subreddit",
            "history",
            "memory",
            "feature",
            "validate",
            "train",
        ]:
            meta["keywords"][key] = small_text.count(key)

    grade_data = load_json(grade_path)
    if grade_data is not None:
        best_score = None
        best_node = None
        for d in iter_dicts(grade_data):
            node_id = d.get("node_id") or d.get("id") or d.get("node")
            candidates = []
            for key in ["score", "best_score", "auc", "roc_auc", "validation_score"]:
                value = d.get(key)
                if isinstance(value, (int, float)):
                    candidates.append(float(value))
            if candidates:
                score = max(candidates)
                if best_score is None or score > best_score:
                    best_score = score
                    best_node = str(node_id) if node_id is not None else None
        meta["best_node"] = best_node
        meta["best_score"] = best_score

    meta["ok"] = node_data is not None or traj_data is not None
    return meta


def node_center(name: str) -> Tuple[float, float]:
    item = TREE[name]
    if item["kind"] == "node":
        return item["cx"], item["cy"]
    return item["x"] + item["w"] / 2, item["y"] + item["h"] / 2


def edge_points(start: str, end: str, pad: float = 0.08) -> Tuple[float, float, float, float]:
    x1, y1 = node_center(start)
    x2, y2 = node_center(end)
    dx = x2 - x1
    dy = y2 - y1
    dist = math.hypot(dx, dy) or 1.0

    def radius(name: str) -> Tuple[float, float]:
        item = TREE[name]
        if item["kind"] == "node":
            r = item["d"] / 2
            return r, r
        return item["w"] / 2, item["h"] / 2

    r1x, r1y = radius(start)
    r2x, r2y = radius(end)
    # A simple directional padding is enough for this schematic; using the
    # smaller projected radius prevents arrows from entering node labels.
    p1 = min(r1x, r1y) + pad
    p2 = min(r2x, r2y) + pad
    return x1 + dx / dist * p1, y1 + dy / dist * p1, x2 - dx / dist * p2, y2 - dy / dist * p2


def setup_text_frame(shape, margin: float = 0.03) -> Any:
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tf


def set_run(run, text: str, size: float, color: str = "ink", bold: bool = False) -> None:
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = ppt_rgb(COLORS[color])


def add_round_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    line: str,
    radius: bool = True,
    line_width: float = 1.0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(COLORS[fill])
    shape.line.color.rgb = ppt_rgb(COLORS[line])
    shape.line.width = Pt(line_width)
    return shape


def add_textbox(slide, x: float, y: float, w: float, h: float, lines: List[Tuple[str, float, str, bool]], align=PP_ALIGN.CENTER) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = setup_text_frame(box, margin=0.01)
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = 0.86
        run = p.add_run()
        set_run(run, text, size, color, bold)


def add_ppt_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: str = "line_dark", width: float = 1.2, dashed: bool = False, arrow: bool = True) -> None:
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = ppt_rgb(COLORS[color])
    connector.line.width = Pt(width)
    if dashed:
        connector.line.dash_style = 4
    if not arrow:
        return
    angle = math.degrees(math.atan2(y2 - y1, x2 - x1))
    size = 0.075 + 0.012 * min(width, 2.0)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x2 - size / 2), Inches(y2 - size / 2), Inches(size), Inches(size))
    tri.fill.solid()
    tri.fill.fore_color.rgb = ppt_rgb(COLORS[color])
    tri.line.color.rgb = ppt_rgb(COLORS[color])
    tri.rotation = angle + 90


def add_node(slide, name: str, label: str, fill: str, line: str, text_color: str = "ink", line_width: float = 1.3, d: Optional[float] = None) -> None:
    item = TREE[name]
    diameter = d or item.get("d", NODE_D)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(item["cx"] - diameter / 2),
        Inches(item["cy"] - diameter / 2),
        Inches(diameter),
        Inches(diameter),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(COLORS[fill])
    shape.line.color.rgb = ppt_rgb(COLORS[line])
    shape.line.width = Pt(line_width)
    tf = setup_text_frame(shape, margin=0.01)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, label, 8.5 if label != "Best" else 8.0, text_color, True)


def add_step_callout(slide, step: int, label: str, x: float, y: float, w: float, h: float, target: str, side: str, accent: str) -> None:
    rect = add_round_rect(slide, x, y, w, h, "callout_fill", "callout_line", radius=True, line_width=0.7)
    setup_text_frame(rect, margin=0.02)
    bar_x = x + 0.10 if side == "left" else x + w - 0.14
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bar_x), Inches(y + 0.08), Inches(0.035), Inches(h - 0.16))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ppt_rgb(COLORS[accent])
    bar.line.color.rgb = ppt_rgb(COLORS[accent])

    label_x = x + 0.24 if side == "left" else x + 0.14
    add_textbox(slide, label_x, y + 0.045, 0.62, h - 0.07, [(f"Step {step}", 7.1, "ink", True)], align=PP_ALIGN.LEFT)
    add_textbox(slide, label_x + 0.72, y + 0.045, w - 1.02, h - 0.07, [(label, 6.95, "muted", False)], align=PP_ALIGN.LEFT)

    tx, ty = node_center(target)
    if side == "left":
        add_ppt_arrow(slide, x + w, y + h / 2, tx - 0.31, ty, color="line", width=0.7, arrow=False)
    else:
        add_ppt_arrow(slide, x, y + h / 2, tx + 0.31, ty, color="line", width=0.7, arrow=False)


def add_memory_box(slide) -> None:
    x, y, w, h = MEMORY_BOX["x"], MEMORY_BOX["y"], MEMORY_BOX["w"], MEMORY_BOX["h"]
    rect = add_round_rect(slide, x, y, w, h, "memory_fill", "memory_line", radius=True, line_width=1.1)
    tf = setup_text_frame(rect, margin=0.11)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(1)
    run = p.add_run()
    set_run(run, "Global Memory", 9.0, "ink", True)
    for item in ["useful hints", "failed attempts", "reusable features", "validation feedback"]:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        p.line_spacing = 0.92
        run = p.add_run()
        set_run(run, f"- {item}", 7.1, "ink", False)


def add_role_box(slide) -> None:
    x, y, w, h = ROLE_BOX["x"], ROLE_BOX["y"], ROLE_BOX["w"], ROLE_BOX["h"]
    rect = add_round_rect(slide, x, y, w, h, "white", "callout_line", radius=True, line_width=0.8)
    setup_text_frame(rect, margin=0.03)
    add_textbox(slide, x + 0.20, y + 0.14, 0.62, 0.22, [("RED", 7.8, "red_line", True)], align=PP_ALIGN.LEFT)
    add_textbox(slide, x + 1.02, y + 0.14, 2.20, 0.22, [("exploration / search / inspection", 6.7, "muted", False)], align=PP_ALIGN.LEFT)
    add_textbox(slide, x + 0.20, y + 0.50, 0.75, 0.22, [("BLACK", 7.8, "ink", True)], align=PP_ALIGN.LEFT)
    add_textbox(slide, x + 1.02, y + 0.45, 2.20, 0.36, [("execution / curation / cleaning", 6.6, "muted", False), ("feature pipeline / validation", 6.6, "muted", False)], align=PP_ALIGN.LEFT)


def build_pptx(path: Path) -> None:
    if Presentation is None:
        raise RuntimeError("python-pptx is not available")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ppt_rgb(COLORS["white"])

    # Faint vertical guide emphasizes that the search process proceeds downward.
    add_ppt_arrow(slide, 6.30, 0.92, 6.30, 6.11, color="faint", width=1.0, dashed=True, arrow=False)

    # Low-contrast sibling branch: enough to signal tree search without clutter.
    add_ppt_arrow(slide, 6.53, 1.18, 7.88, 1.23, color="faint", width=1.0, arrow=False)
    for cx, cy, fill, line, label, txt in [
        (8.16, 1.23, "side_red", "red_line", "R", "red_line"),
        (8.78, 1.23, "side_black", "line_dark", "B", "muted"),
    ]:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - SMALL_D / 2), Inches(cy - SMALL_D / 2), Inches(SMALL_D), Inches(SMALL_D))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_rgb(COLORS[fill])
        shape.line.color.rgb = ppt_rgb(COLORS[line])
        shape.line.width = Pt(0.8)
        tf = setup_text_frame(shape, margin=0.005)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        set_run(run, label, 5.8, txt, True)
    add_ppt_arrow(slide, 8.30, 1.23, 8.64, 1.23, color="faint", width=0.8, arrow=True)
    add_textbox(slide, 8.03, 1.43, 0.95, 0.18, [("side branches", 5.8, "muted", False)], align=PP_ALIGN.CENTER)

    # Task, cycle marker, and output.
    task = TREE["task"]
    add_round_rect(slide, task["x"], task["y"], task["w"], task["h"], "task_fill", "task_line", radius=True, line_width=1.2)
    add_textbox(slide, task["x"] + 0.13, task["y"] + 0.09, task["w"] - 0.26, 0.20, [("Task Input", 8.3, "ink", True)])
    add_textbox(slide, task["x"] + 0.13, task["y"] + 0.32, task["w"] - 0.26, 0.16, [("random-acts-of-pizza", 6.8, "muted", False)])
    add_textbox(slide, task["x"] + 0.13, task["y"] + 0.45, task["w"] - 0.26, 0.12, [("request text + metadata", 5.8, "muted", False)])

    cycle = TREE["cycle"]
    add_round_rect(slide, cycle["x"], cycle["y"], cycle["w"], cycle["h"], "white", "line", radius=True, line_width=0.9)
    add_textbox(slide, cycle["x"] + 0.10, cycle["y"] + 0.05, cycle["w"] - 0.20, 0.14, [("...", 14.0, "line_dark", True)])
    add_textbox(slide, cycle["x"] + 0.10, cycle["y"] + 0.24, cycle["w"] - 0.20, 0.11, [("repeated RED/BLACK cycles", 5.9, "muted", False)])

    output = TREE["output"]
    add_round_rect(slide, output["x"], output["y"], output["w"], output["h"], "output_fill", "output_line", radius=True, line_width=1.2)
    add_textbox(slide, output["x"] + 0.13, output["y"] + 0.13, output["w"] - 0.26, 0.20, [("Best Submission", 8.4, "ink", True)])
    add_textbox(slide, output["x"] + 0.13, output["y"] + 0.39, output["w"] - 0.26, 0.14, [("final prediction file", 6.4, "muted", False)])

    # Main nodes.
    add_node(slide, "init", "Init", "init_fill", "init_line", "ink", line_width=1.2)
    for name in ["R1", "R2", "R3", "Rk"]:
        add_node(slide, name, name, "red_fill", "red_line", "ink", line_width=1.2)
    for name in ["B1", "B2", "B3", "Bk"]:
        add_node(slide, name, name, "black_fill", "black_line", "white", line_width=1.2)
    add_node(slide, "Best", "Best", "black_fill", "gold", "white", line_width=2.2, d=BEST_D)
    star = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(6.58), Inches(5.08), Inches(0.16), Inches(0.16))
    star.fill.solid()
    star.fill.fore_color.rgb = ppt_rgb(COLORS["gold"])
    star.line.color.rgb = ppt_rgb(COLORS["gold"])

    # Main path after nodes so arrowheads are visible but do not cover labels too much.
    for start, end in MAIN_EDGES:
        x1, y1, x2, y2 = edge_points(start, end)
        add_ppt_arrow(slide, x1, y1, x2, y2, color="line_dark", width=1.25, arrow=True)

    # Memory and role definitions.
    add_role_box(slide)
    add_memory_box(slide)
    add_ppt_arrow(slide, 7.30, 3.35, 9.10, 4.84, color="memory_line", width=1.0, arrow=True)
    add_textbox(slide, 7.72, 4.05, 1.02, 0.18, [("Add / Update", 6.2, "memory_text", False)], align=PP_ALIGN.CENTER)
    add_ppt_arrow(slide, 9.10, 5.34, 5.86, 4.68, color="memory_line", width=1.0, arrow=True)
    add_textbox(slide, 7.25, 5.05, 1.00, 0.18, [("Read / Reuse", 6.2, "memory_text", False)], align=PP_ALIGN.CENTER)

    # Step callouts are drawn last so their labels stay crisp.
    for item in STEP_CALLOUTS:
        add_step_callout(slide, *item)

    prs.save(path)


def export_with_libreoffice(pptx_path: Path, out_dir: Path) -> bool:
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        print("Warning: LibreOffice/soffice not found; using matplotlib fallback for PDF/PNG.")
        return False
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
        return PDF_PATH.exists()
    except Exception as exc:
        print(f"Warning: LibreOffice export failed ({exc}); using matplotlib fallback for PDF/PNG.")
        return False


def mpl_color(name: str) -> str:
    return "#" + COLORS[name]


def mpl_round_rect(ax, x, y, w, h, fill, line, lw=1.0, radius=0.08, z=2):
    patch = FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle=f"round,pad=0.02,rounding_size={radius}",
        facecolor=mpl_color(fill),
        edgecolor=mpl_color(line),
        linewidth=lw,
        zorder=z,
    )
    ax.add_patch(patch)
    return patch


def mpl_text(ax, x, y, text, size=7, color="ink", weight="normal", ha="center", va="center", z=5):
    ax.text(
        x,
        y,
        text,
        fontsize=size,
        color=mpl_color(color),
        fontweight=weight,
        ha=ha,
        va=va,
        family=MPL_FONT,
        linespacing=0.88,
        zorder=z,
    )


def mpl_arrow(ax, x1, y1, x2, y2, color="line_dark", lw=1.1, arrow=True, dashed=False, z=1, curve=0.0):
    style = "-|>" if arrow else "-"
    patch = FancyArrowPatch(
        (x1, y1),
        (x2, y2),
        arrowstyle=style,
        mutation_scale=8,
        linewidth=lw,
        color=mpl_color(color),
        linestyle=(0, (4, 3)) if dashed else "solid",
        connectionstyle=f"arc3,rad={curve}",
        zorder=z,
    )
    ax.add_patch(patch)


def mpl_node(ax, name, label, fill, line, text_color="ink", lw=1.2, d=None, z=4):
    item = TREE[name]
    diameter = d or item.get("d", NODE_D)
    circle = Circle((item["cx"], item["cy"]), diameter / 2, facecolor=mpl_color(fill), edgecolor=mpl_color(line), linewidth=lw, zorder=z)
    ax.add_patch(circle)
    mpl_text(ax, item["cx"], item["cy"], label, size=8.0 if label != "Best" else 7.6, color=text_color, weight="bold", z=z + 1)


def mpl_step(ax, step, label, x, y, w, h, target, side, accent):
    mpl_round_rect(ax, x, y, w, h, "callout_fill", "callout_line", lw=0.7, radius=0.06, z=4)
    bx = x + 0.10 if side == "left" else x + w - 0.14
    ax.add_patch(FancyBboxPatch((bx, y + 0.08), 0.035, h - 0.16, boxstyle="square,pad=0", facecolor=mpl_color(accent), edgecolor=mpl_color(accent), linewidth=0, zorder=5))
    label_x = x + 0.24 if side == "left" else x + 0.14
    mpl_text(ax, label_x, y + h / 2, f"Step {step}", size=6.7, weight="bold", ha="left", color="ink", z=6)
    mpl_text(ax, label_x + 0.72, y + h / 2, label, size=6.55, ha="left", color="muted", z=6)
    tx, ty = node_center(target)
    if side == "left":
        mpl_arrow(ax, x + w, y + h / 2, tx - 0.31, ty, color="line", lw=0.7, arrow=False, z=1)
    else:
        mpl_arrow(ax, x, y + h / 2, tx + 0.31, ty, color="line", lw=0.7, arrow=False, z=1)


def build_matplotlib_outputs(pdf_path: Path, png_path: Path) -> None:
    plt.rcParams["pdf.fonttype"] = 42
    plt.rcParams["ps.fonttype"] = 42
    plt.rcParams["font.family"] = MPL_FONT
    fig = plt.figure(figsize=(SLIDE_W, SLIDE_H), dpi=160)
    fig.patch.set_facecolor("white")
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, SLIDE_W)
    ax.set_ylim(SLIDE_H, 0)
    ax.axis("off")

    mpl_arrow(ax, 6.30, 0.92, 6.30, 6.11, color="faint", lw=1.0, arrow=False, dashed=True, z=0)

    # Side branch.
    mpl_arrow(ax, 6.53, 1.18, 7.88, 1.23, color="faint", lw=1.0, arrow=False, z=0)
    for cx, cy, fill, line, label, text_col in [
        (8.16, 1.23, "side_red", "red_line", "R", "red_line"),
        (8.78, 1.23, "side_black", "line_dark", "B", "muted"),
    ]:
        ax.add_patch(Circle((cx, cy), SMALL_D / 2, facecolor=mpl_color(fill), edgecolor=mpl_color(line), linewidth=0.8, zorder=1))
        mpl_text(ax, cx, cy, label, size=5.5, color=text_col, weight="bold", z=2)
    mpl_arrow(ax, 8.30, 1.23, 8.64, 1.23, color="faint", lw=0.8, arrow=True, z=0)
    mpl_text(ax, 8.50, 1.54, "side branches", size=5.6, color="muted", z=1)

    # Main boxes.
    task = TREE["task"]
    mpl_round_rect(ax, task["x"], task["y"], task["w"], task["h"], "task_fill", "task_line", lw=1.2, radius=0.09, z=3)
    mpl_text(ax, task["x"] + task["w"] / 2, task["y"] + 0.20, "Task Input", size=8.0, weight="bold")
    mpl_text(ax, task["x"] + task["w"] / 2, task["y"] + 0.40, "random-acts-of-pizza", size=6.4, color="muted")
    mpl_text(ax, task["x"] + task["w"] / 2, task["y"] + 0.53, "request text + metadata", size=5.5, color="muted")

    cycle = TREE["cycle"]
    mpl_round_rect(ax, cycle["x"], cycle["y"], cycle["w"], cycle["h"], "white", "line", lw=0.9, radius=0.07, z=3)
    mpl_text(ax, cycle["x"] + cycle["w"] / 2, cycle["y"] + 0.14, "...", size=13.0, color="line_dark", weight="bold")
    mpl_text(ax, cycle["x"] + cycle["w"] / 2, cycle["y"] + 0.31, "repeated RED/BLACK cycles", size=5.6, color="muted")

    output = TREE["output"]
    mpl_round_rect(ax, output["x"], output["y"], output["w"], output["h"], "output_fill", "output_line", lw=1.2, radius=0.08, z=3)
    mpl_text(ax, output["x"] + output["w"] / 2, output["y"] + 0.27, "Best Submission", size=8.0, weight="bold")
    mpl_text(ax, output["x"] + output["w"] / 2, output["y"] + 0.49, "final prediction file", size=6.0, color="muted")

    # Nodes and main path.
    mpl_node(ax, "init", "Init", "init_fill", "init_line", "ink", lw=1.2)
    for name in ["R1", "R2", "R3", "Rk"]:
        mpl_node(ax, name, name, "red_fill", "red_line", "ink", lw=1.2)
    for name in ["B1", "B2", "B3", "Bk"]:
        mpl_node(ax, name, name, "black_fill", "black_line", "white", lw=1.2)
    mpl_node(ax, "Best", "Best", "black_fill", "gold", "white", lw=2.2, d=BEST_D)
    star = Polygon([[6.66, 5.08], [6.69, 5.17], [6.79, 5.17], [6.71, 5.22], [6.74, 5.31], [6.66, 5.25], [6.58, 5.31], [6.61, 5.22], [6.53, 5.17], [6.63, 5.17]], closed=True, facecolor=mpl_color("gold"), edgecolor=mpl_color("gold"), zorder=5)
    ax.add_patch(star)

    for start, end in MAIN_EDGES:
        x1, y1, x2, y2 = edge_points(start, end)
        mpl_arrow(ax, x1, y1, x2, y2, color="line_dark", lw=1.15, arrow=True, z=2)

    # Memory and role key.
    x, y, w, h = ROLE_BOX["x"], ROLE_BOX["y"], ROLE_BOX["w"], ROLE_BOX["h"]
    mpl_round_rect(ax, x, y, w, h, "white", "callout_line", lw=0.8, radius=0.08, z=3)
    mpl_text(ax, x + 0.22, y + 0.27, "RED", size=7.4, color="red_line", weight="bold", ha="left")
    mpl_text(ax, x + 1.02, y + 0.27, "exploration / search / inspection", size=6.25, color="muted", ha="left")
    mpl_text(ax, x + 0.22, y + 0.62, "BLACK", size=7.4, color="ink", weight="bold", ha="left")
    mpl_text(ax, x + 1.02, y + 0.55, "execution / curation / cleaning", size=6.1, color="muted", ha="left")
    mpl_text(ax, x + 1.02, y + 0.73, "feature pipeline / validation", size=6.1, color="muted", ha="left")

    x, y, w, h = MEMORY_BOX["x"], MEMORY_BOX["y"], MEMORY_BOX["w"], MEMORY_BOX["h"]
    mpl_round_rect(ax, x, y, w, h, "memory_fill", "memory_line", lw=1.1, radius=0.08, z=3)
    mpl_text(ax, x + 0.22, y + 0.30, "Global Memory", size=8.5, color="ink", weight="bold", ha="left")
    for idx, item in enumerate(["useful hints", "failed attempts", "reusable features", "validation feedback"]):
        mpl_text(ax, x + 0.28, y + 0.55 + idx * 0.20, f"- {item}", size=6.65, color="ink", ha="left")

    mpl_arrow(ax, 7.30, 3.35, 9.10, 4.84, color="memory_line", lw=1.0, arrow=True, z=2, curve=0.04)
    mpl_text(ax, 8.23, 4.13, "Add / Update", size=5.9, color="memory_text")
    mpl_arrow(ax, 9.10, 5.34, 5.86, 4.68, color="memory_line", lw=1.0, arrow=True, z=2, curve=0.05)
    mpl_text(ax, 7.75, 5.13, "Read / Reuse", size=5.9, color="memory_text")

    for item in STEP_CALLOUTS:
        mpl_step(ax, *item)

    fig.savefig(pdf_path, facecolor="white")
    fig.savefig(png_path, dpi=350, facecolor="white")
    plt.close(fig)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    meta = parse_run_logs()
    if not meta["ok"]:
        print("Warning: could not parse run logs; using fallback semi-concrete random-acts-of-pizza labels.")
    elif meta["warnings"]:
        print("Warning: parsed partial run logs; using schematic labels with fallback where needed.")
        for warning in meta["warnings"][:2]:
            print(f"  - {warning}")
    else:
        best = meta.get("best_node")
        score = meta.get("best_score")
        score_text = f", best score {score:.5f}" if isinstance(score, float) else ""
        best_text = f", best node {best[:8]}" if isinstance(best, str) and best else ""
        print(
            "Parsed run logs: "
            f"{meta.get('node_count', 0)} UCT nodes, "
            f"{meta.get('trajectory_records', 0)} trajectory records"
            f"{best_text}{score_text}."
        )

    build_pptx(PPTX_PATH)
    exported = export_with_libreoffice(PPTX_PATH, OUT_DIR)
    if not exported or not PDF_PATH.exists() or not PNG_PATH.exists():
        build_matplotlib_outputs(PDF_PATH, PNG_PATH)

    print("Generated files:")
    print(PPTX_PATH)
    print(PDF_PATH)
    print(PNG_PATH)


if __name__ == "__main__":
    main()
